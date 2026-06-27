"""
faire_pptx.py
=============
Genere la presentation PowerPoint du projet (presentation.pptx).
Palette navale "Midnight Executive" : navy / bleu glace / blanc, accent or.
Les 4 figures du projet sont integrees.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ICI = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(ICI, "..", "figures")

# --- Palette ---
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY2 = RGBColor(0x2A, 0x35, 0x80)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF2, 0xA9, 0x00)
DARK = RGBColor(0x21, 0x29, 0x34)
MUTED = RGBColor(0x6B, 0x72, 0x80)

TITRE_FONT = "Cambria"
CORPS_FONT = "Calibri"

# Aspect (largeur/hauteur) des figures pour caler la mise en page
ASPECT = {
    "1_plan_des_formes.png": 7 / 6,
    "2_courbe_des_aires.png": 7 / 4,
    "3_vue_3d.png": 10 / 5,
    "4_courbe_gz.png": 7 / 5,
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def texte(s, l, t, w, h, txt, size, color, bold=False, font=CORPS_FONT,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lignes = txt.split("\n")
    for i, ligne in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ligne
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def carte(s, l, t, w, h, fill, ligne=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if ligne is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = ligne
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def cercle(s, l, t, d, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def image(s, nom, l, t, w):
    """Place une figure par sa largeur ; hauteur calculee via l'aspect."""
    h = w / ASPECT[nom]
    s.shapes.add_picture(os.path.join(FIG, nom), Inches(l), Inches(t),
                         width=Inches(w), height=Inches(h))
    return h


# =====================================================================
# Slide 1 — Titre (fond navy)
# =====================================================================
s = slide(NAVY)
texte(s, 0.7, 1.6, 7.2, 2.2,
      "Carène de frégate", 54, WHITE, bold=True, font=TITRE_FONT)
texte(s, 0.72, 3.1, 7.0, 1.2,
      "Hydrostatique & stabilité\nOutil Python d'architecture navale", 22, ICE,
      font=CORPS_FONT, spacing=1.1)
texte(s, 0.72, 5.4, 7.2, 0.5,
      "Ayoub — Licence Physique CUPGE, UBS Lorient", 16, WHITE, bold=True)
texte(s, 0.72, 5.95, 7.2, 0.5,
      "Projet personnel — candidature ENSTA / Naval Group", 13, ICE, italic=True)
# carte image 3D
carte(s, 8.25, 1.5, 4.55, 3.0, WHITE)
image(s, "3_vue_3d.png", 8.45, 2.35, 4.15)

# =====================================================================
# Slide 2 — Objectif & démarche (3 piliers)
# =====================================================================
s = slide(WHITE)
texte(s, 0.7, 0.55, 12, 1.0, "Objectif & démarche", 40, NAVY, bold=True, font=TITRE_FONT)
texte(s, 0.72, 1.55, 11.8, 1.0,
      "Concevoir une carène de frégate et calculer sa flottabilité et sa stabilité "
      "par le code. Le projet fait converger trois compétences en un seul livrable.",
      18, DARK, spacing=1.1)

piliers = [
    ("Physique", "Archimède, moments,\nmétacentre, stabilité", NAVY),
    ("Géométrie", "Carène paramétrique\ngénérée par formule", GOLD),
    ("Code", "Intégration de Simpson,\ncalcul numérique", NAVY2),
]
x = 0.9
for nom, desc, coul in piliers:
    cercle(s, x + 1.0, 3.1, 1.1, coul)
    texte(s, x, 4.35, 3.3, 0.6, nom, 22, NAVY, bold=True, align=PP_ALIGN.CENTER, font=TITRE_FONT)
    texte(s, x, 5.0, 3.3, 1.2, desc, 15, MUTED, align=PP_ALIGN.CENTER, spacing=1.05)
    x += 4.0

# =====================================================================
# Slide 3 — La carène paramétrique
# =====================================================================
s = slide(WHITE)
texte(s, 0.7, 0.55, 12, 1.0, "La carène paramétrique", 40, NAVY, bold=True, font=TITRE_FONT)
texte(s, 0.72, 1.7, 6.2, 0.6, "Une coque décrite par une formule", 20, GOLD, bold=True)
texte(s, 0.72, 2.45, 6.2, 0.8,
      "b(x, z) = (B/2) · P(x) · V(z)", 20, DARK, bold=True, font="Courier New")
texte(s, 0.72, 3.4, 6.2, 3.3,
      "P(x) : affinement de l'avant vers l'arrière\n"
      "V(z) : forme en U de la quille au pont\n\n"
      "Quelques paramètres suffisent :\n"
      "L = 100 m   B = 14 m\n"
      "T = 4.5 m   D = 9 m\n\n"
      "Changer un paramètre régénère toute la coque.",
      17, DARK, spacing=1.15)
carte(s, 7.35, 1.7, 5.5, 5.2, ICE)
image(s, "1_plan_des_formes.png", 7.75, 2.0, 4.7)

# =====================================================================
# Slide 4 — Méthode : Simpson + courbe des aires
# =====================================================================
s = slide(WHITE)
texte(s, 0.7, 0.55, 12, 1.0, "Méthode de calcul", 40, NAVY, bold=True, font=TITRE_FONT)
texte(s, 0.72, 1.7, 6.0, 0.6, "L'intégration de Simpson", 20, GOLD, bold=True)
texte(s, 0.72, 2.4, 6.0, 4.2,
      "On calcule volumes et aires en intégrant la géométrie le long de la coque.\n\n"
      "La méthode de Simpson approche la courbe par des paraboles : c'est la "
      "méthode historique des architectes navals, codée à la main ici.\n\n"
      "La courbe des aires de section A(x) (ci-contre) est la première brique : "
      "son intégrale donne le volume immergé, donc le déplacement.",
      17, DARK, spacing=1.15)
carte(s, 7.15, 2.0, 5.7, 3.7, ICE)
image(s, "2_courbe_des_aires.png", 7.5, 2.45, 5.0)

# =====================================================================
# Slide 5 — Résultats hydrostatiques (stat grid)
# =====================================================================
s = slide(NAVY)
texte(s, 0.7, 0.55, 12, 1.0, "Résultats hydrostatiques", 40, WHITE, bold=True, font=TITRE_FONT)
texte(s, 0.72, 1.55, 12, 0.6, "Navire droit, au tirant d'eau de projet", 17, ICE, italic=True)

stats = [
    ("2 870 t", "Déplacement"),
    ("0.444", "Coefficient de bloc Cb"),
    ("1.15 m", "Hauteur métacentrique GMt"),
    ("6.55 m", "Métacentre KMt"),
]
cw, ch, gap = 5.7, 1.95, 0.5
x0, y0 = 0.85, 2.55
for i, (val, lab) in enumerate(stats):
    cx = x0 + (i % 2) * (cw + gap)
    cy = y0 + (i // 2) * (ch + 0.4)
    carte(s, cx, cy, cw, ch, NAVY2)
    texte(s, cx + 0.4, cy + 0.25, cw - 0.8, 1.0, val, 44, GOLD, bold=True, font=TITRE_FONT)
    texte(s, cx + 0.4, cy + 1.25, cw - 0.8, 0.6, lab, 15, ICE)

# =====================================================================
# Slide 6 — Stabilité (courbe GZ)
# =====================================================================
s = slide(WHITE)
texte(s, 0.7, 0.55, 12, 1.0, "Stabilité : la courbe GZ", 40, NAVY, bold=True, font=TITRE_FONT)
hh = image(s, "4_courbe_gz.png", 6.6, 1.75, 6.3)
# indicateurs a gauche
ind = [
    ("GZ max", "0.59 m à 45°"),
    ("Limite de stabilité", "83°"),
    ("GM (pente initiale)", "1.13 m"),
]
y = 2.0
for lab, val in ind:
    carte(s, 0.8, y, 5.2, 1.2, ICE)
    texte(s, 1.1, y + 0.18, 4.7, 0.5, lab, 15, MUTED, bold=True)
    texte(s, 1.1, y + 0.55, 4.7, 0.6, val, 24, NAVY, bold=True, font=TITRE_FONT)
    y += 1.5

# =====================================================================
# Slide 7 — Validation (crédibilité)
# =====================================================================
s = slide(WHITE)
texte(s, 0.7, 0.55, 12, 1.0, "Validation du calcul", 40, NAVY, bold=True, font=TITRE_FONT)
texte(s, 0.72, 1.6, 11.8, 1.0,
      "Le GM est calculé de deux façons indépendantes. Elles concordent : "
      "le résultat est fiable.", 18, DARK, spacing=1.1)
duo = [
    ("1.146 m", "GM par l'hydrostatique", NAVY),
    ("1.131 m", "GM par la pente de la courbe GZ", GOLD),
]
x = 1.2
for val, lab, coul in duo:
    carte(s, x, 3.0, 5.0, 2.4, WHITE, ligne=coul)
    texte(s, x + 0.4, 3.45, 4.2, 1.0, val, 50, coul, bold=True, font=TITRE_FONT)
    texte(s, x + 0.4, 4.55, 4.2, 0.7, lab, 16, DARK)
    x += 5.7
texte(s, 1.2, 5.8, 11, 0.8,
      "Écart inférieur à 1.5 %. La formule des murailles droites confirme aussi "
      "la courbe aux petits angles.", 16, MUTED, italic=True)

# =====================================================================
# Slide 8 — Exports CAO & reproductibilité
# =====================================================================
s = slide(WHITE)
texte(s, 0.7, 0.55, 12, 1.0, "Exports vers la CAO", 40, NAVY, bold=True, font=TITRE_FONT)
texte(s, 0.72, 1.6, 11.8, 0.7,
      "La géométrie est exportée pour être finalisée dans Fusion 360.", 18, DARK)
cards = [
    ("STL", "Maillage 3D de la coque\n→ Insertion > Maillage"),
    ("DXF", "Sections (couples)\n→ à lisser par loft"),
    ("CSV", "Table des cotes\n(demi-largeurs)"),
]
x = 0.9
for nom, desc in cards:
    carte(s, x, 2.8, 3.7, 2.8, ICE)
    texte(s, x, 3.1, 3.7, 0.8, nom, 30, NAVY, bold=True, align=PP_ALIGN.CENTER, font=TITRE_FONT)
    texte(s, x + 0.4, 4.1, 2.9, 1.3, desc, 15, DARK, align=PP_ALIGN.CENTER, spacing=1.1)
    x += 4.0
texte(s, 0.9, 6.0, 11.5, 0.8,
      "Tout est régénérable par une seule commande : python main.py", 16, MUTED, italic=True)

# =====================================================================
# Slide 9 — Conclusion / suite (fond navy)
# =====================================================================
s = slide(NAVY)
texte(s, 0.7, 1.0, 12, 1.0, "Bilan & suite", 42, WHITE, bold=True, font=TITRE_FONT)
texte(s, 0.72, 2.2, 11.8, 2.0,
      "Un outil complet : géométrie paramétrique, hydrostatique, stabilité GZ, "
      "exports CAO. Calcul validé et entièrement reproductible.",
      20, ICE, spacing=1.2)
texte(s, 0.72, 4.0, 11.8, 0.6, "Phase 2 (2027)", 22, GOLD, bold=True, font=TITRE_FONT)
texte(s, 0.72, 4.75, 11.8, 2.0,
      "• Évasement de coque et formes plus réalistes\n"
      "• Couplage gîte / assiette\n"
      "• Résistance à l'avancement et puissance moteur\n"
      "• Étude d'un navire de soutien offshore (OSV)",
      18, WHITE, spacing=1.25)

sortie = os.path.join(ICI, "presentation.pptx")
prs.save(sortie)
print("PowerPoint genere :", sortie, "(", len(prs.slides.__iter__.__self__._sldIdLst), "slides )")
